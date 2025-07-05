import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from tkcalendar import DateEntry
import openpyxl
from openpyxl import Workbook
import os
import re
import pyperclip
from datetime import datetime
import matplotlib.pyplot as plt
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
import subprocess

from pptx import Presentation
from pptx.util import Pt
import sys

# Función para ubicar recursos en un ejecutable PyInstaller
def recurso_relativo(ruta_relativa):
    """Obtiene la ruta absoluta del recurso, compatible con PyInstaller."""
    try:
        # PyInstaller crea una carpeta temporal en _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
    
    return os.path.join(base_path, ruta_relativa)

class BitacoraCTApp:
    def __init__(self, root):
        self.root = root
        self.version = "1.3.2"  # Versión actualizada
        self.ruta_base_datos = r"\\mexhome03\Data\Prototype Engineering\Public\MC Front End\Mold\PROCESOS MOLDEO\ramon\01- Documentos\Bitacora CT\Bitacora CT.xlsx"
        self.id_a_modificar = None

        # Definir las listas para los Combobox
        self.combo_list_Equipo = ["Towa20", "Towa21", "Towa22", "Towa23", "Towa24", "Towa25",
                                 "Towa26", "Towa27", "Towa28", "Towa29", "Towa30", "Towa31",
                                 "Towa32", "Towa33", "Otro"]
        self.combo_list_Turno = ["A", "B", "C", "D"]
        self.combo_list_RCA = ["Si", "No", "Pendiente"]
        self.combo_list_Estatus = ["Quitar Hold y Mover", "Descontar Quitar Hold y Mover", "Para NCMR","Pendiente"]
        
        self.configurar_interfaz()
        self.crear_widgets()
        self.actualizar_treeview()
        self.centrar_ventana()

    def configurar_interfaz(self):
        self.root.geometry("1200x650")
        self.root.title(f"Bitácora CT Hold - v{self.version}")
        
        # Configuración del tema
        style = ttk.Style(self.root)
        self.root.option_add("*tearOff", False)
        self.root.tk.call("source", recurso_relativo("forest-light.tcl"))
        self.root.tk.call("source", recurso_relativo("forest-dark.tcl"))
        style.theme_use("forest-dark")

        # Frame principal
        self.frame = ttk.Frame(self.root)
        self.frame.pack(fill="both", expand=True, padx=10, pady=10)

    def crear_widgets(self):
        self.crear_formulario()
        self.crear_treeview()
        self.crear_botones_adicionales()

    def crear_formulario(self):
        # Frame principal del formulario
        widgets_frame = ttk.LabelFrame(self.frame, text="Insertar Evento")
        widgets_frame.grid(row=0, column=0, padx=10, pady=10, sticky="nsew")
        
        # Configurar grid para mejor distribución
        widgets_frame.columnconfigure(0, weight=1)
        widgets_frame.columnconfigure(1, weight=1)
        
        # Frame para campos del formulario
        campos_frame = ttk.Frame(widgets_frame)
        campos_frame.grid(row=0, column=0, columnspan=2, sticky="nsew", pady=(0, 10))
        
        # Lote
        self.lot_entry = ttk.Entry(campos_frame)
        self.lot_entry.insert(0, "Lote")
        self.lot_entry.bind("<FocusIn>", lambda e: self.lot_entry.delete(0, "end") if self.lot_entry.get() == "Lote" else None)
        self.lot_entry.grid(row=0, column=0, padx=5, pady=5, sticky="ew", columnspan=2)

        # Part ID
        self.part_id_entry = ttk.Entry(campos_frame)
        self.part_id_entry.insert(0, "Part ID")
        self.part_id_entry.bind("<FocusIn>", lambda e: self.part_id_entry.delete(0, "end") if self.part_id_entry.get() == "Part ID" else None)
        self.part_id_entry.grid(row=1, column=0, padx=5, pady=5, sticky="ew", columnspan=2)

        # Equipo (Combobox)
        self.equipo_combo = ttk.Combobox(campos_frame, values=self.combo_list_Equipo)
        self.equipo_combo.set("Equipo")
        self.equipo_combo.grid(row=2, column=0, padx=5, pady=5, sticky="ew", columnspan=2)

        # CT y CP
        ct_cp_frame = ttk.Frame(campos_frame)
        ct_cp_frame.grid(row=3, column=0, columnspan=2, sticky="ew", pady=5)
        ct_cp_frame.columnconfigure(0, weight=1)
        ct_cp_frame.columnconfigure(1, weight=1)

        self.ct_entry = ttk.Entry(ct_cp_frame, justify="center")
        self.ct_entry.insert(0, "CT")
        self.ct_entry.bind("<FocusIn>", lambda e: self.ct_entry.delete(0, "end") if self.ct_entry.get() == "CT" else None)
        self.ct_entry.grid(row=0, column=0, sticky="ew", padx=(0,2))

        self.cp_entry = ttk.Entry(ct_cp_frame, justify="center")
        self.cp_entry.insert(0, "CP")
        self.cp_entry.bind("<FocusIn>", lambda e: self.cp_entry.delete(0, "end") if self.cp_entry.get() == "CP" else None)
        self.cp_entry.grid(row=0, column=1, sticky="ew", padx=(2,0))

        # Comentario
        ttk.Label(campos_frame, text="Comentario:").grid(row=4, column=0, padx=5, sticky="w", columnspan=2)
        self.cometario_entry = tk.Text(campos_frame, height=6, width=30, wrap="word")
        self.cometario_entry.grid(row=5, column=0, columnspan=2, sticky="ew", pady=5, padx=5)

        # PCB
        self.pcb_entry = ttk.Entry(campos_frame, justify="center")
        self.pcb_entry.insert(0, "PCB")
        self.pcb_entry.bind("<FocusIn>", lambda e: self.pcb_entry.delete(0, "end") if self.pcb_entry.get() == "PCB" else None)
        self.pcb_entry.grid(row=6, column=0, sticky="ew", pady=5, columnspan=2)

        # Fecha y Hora
        fecha_hora_frame = ttk.Frame(campos_frame)
        fecha_hora_frame.grid(row=7, column=0, columnspan=2, sticky="ew", pady=5)
        fecha_hora_frame.columnconfigure(0, weight=1)
        fecha_hora_frame.columnconfigure(1, weight=1)

        ttk.Label(fecha_hora_frame, text="Fecha:").grid(row=0, column=0, padx=(5,2), sticky="w")
        ttk.Label(fecha_hora_frame, text="Hora (24h):").grid(row=0, column=1, padx=(2,5), sticky="w")
        
        self.calendario = DateEntry(fecha_hora_frame, date_pattern="mm/dd/yyyy")
        self.calendario.grid(row=1, column=0, padx=(5,2), sticky="ew")
        
        self.hora_entry = ttk.Entry(fecha_hora_frame)
        self.hora_entry.grid(row=1, column=1, padx=(2,5), sticky="ew")

        # Turno y RCA
        turno_rca_frame = ttk.Frame(campos_frame)
        turno_rca_frame.grid(row=8, column=0, columnspan=2, sticky="ew", pady=5)
        turno_rca_frame.columnconfigure(0, weight=1)
        turno_rca_frame.columnconfigure(1, weight=1)

        self.turno_combo = ttk.Combobox(turno_rca_frame, values=self.combo_list_Turno)
        self.turno_combo.set("Turno")
        self.turno_combo.grid(row=0, column=0, padx=(5,2), sticky="ew")

        self.rca_combo = ttk.Combobox(turno_rca_frame, values=self.combo_list_RCA)
        self.rca_combo.set("RCA")
        self.rca_combo.grid(row=0, column=1, padx=(2,5), sticky="ew")

        # Estatus
        self.estatus_combo = ttk.Combobox(campos_frame, values=self.combo_list_Estatus)
        self.estatus_combo.set("Estatus")
        self.estatus_combo.grid(row=9, column=0, columnspan=2, sticky="ew", pady=5)

        # Separador visual
        ttk.Separator(widgets_frame).grid(row=1, column=0, columnspan=2, sticky="ew", pady=5)

        # Frame para botones principales (ahora en la parte inferior)
        self.button_frame = ttk.Frame(widgets_frame)
        self.button_frame.grid(row=2, column=0, columnspan=2, sticky="ew", pady=(0, 5))

        # Botones principales
        self.insert_button = ttk.Button(self.button_frame, text="Insertar Evento", command=self.guardar_datos, style="Accent.TButton")
        self.insert_button.pack(side="left", padx=2, expand=True)

        ttk.Button(self.button_frame, text="Limpiar", command=self.limpiar_campos).pack(side="left", padx=2, expand=True)
        ttk.Button(self.button_frame, text="Buscar", command=self.buscar_por_lote).pack(side="left", padx=2, expand=True)

        # Frame para botones de gráficas
        self.graph_button_frame = ttk.Frame(widgets_frame)
        self.graph_button_frame.grid(row=3, column=0, columnspan=2, sticky="ew")

        # Botones de gráficas
        ttk.Button(self.graph_button_frame, text="Gráfica Turnos", command=self.generar_grafica_turnos, style="Accent.TButton").pack(side="left", padx=2, expand=True)
        ttk.Button(self.graph_button_frame, text="Gráfica Equipos", command=self.generar_grafica_equipos, style="Accent.TButton").pack(side="left", padx=2, expand=True)

        # Botón de guardar cambios (oculto inicialmente)
        self.guardar_button = ttk.Button(self.button_frame, text="Guardar Cambios", command=self.guardar_cambios, style="Accent.TButton")
        self.guardar_button.pack(side="left", padx=2, expand=True)
        self.guardar_button.pack_forget()

    def crear_treeview(self):
        # Frame para el Treeview
        treeFrame = ttk.LabelFrame(self.frame, text="Registros")
        treeFrame.grid(row=0, column=1, padx=10, pady=10, sticky="nsew")  

        # Scrollbar para el Treeview
        treescroll = ttk.Scrollbar(treeFrame)
        treescroll.pack(side="right", fill="y")

        # Configurar el Treeview
        cols = ("Lote", "Part ID", "Equipo", "Comentario", "Estatus")
        self.treeview = ttk.Treeview(
            treeFrame,
            show="headings",
            yscrollcommand=treescroll.set,
            columns=cols,
            height=12,
            selectmode="browse"
        )

        # Configurar columnas
        self.treeview.column("Lote", anchor="w", width=80)
        self.treeview.column("Part ID", anchor="w", width=80)  
        self.treeview.column("Equipo", anchor="w", width=80)
        self.treeview.column("Comentario", anchor="w", width=200)
        self.treeview.column("Estatus", anchor="w", width=200)

        # Configurar encabezados
        for col in cols:
            self.treeview.heading(col, text=col)

        self.treeview.pack(fill="both", expand=True)
        treescroll.config(command=self.treeview.yview)

    def crear_botones_adicionales(self):
        # Frame para botones del treeview
        botones_tree_frame = ttk.Frame(self.treeview.master)
        botones_tree_frame.pack(fill="x", pady=5)

        # Botón para copiar selección
        ttk.Button(
            botones_tree_frame,
            text="Copiar Selección",
            command=self.copiar_seleccion,
            style="Accent.TButton"
        ).pack(side="left", padx=2, expand=True)

        # Botón para modificar evento
        ttk.Button(
            botones_tree_frame,
            text="Modificar",
            command=self.modificar_evento,
            style="Accent.TButton"
        ).pack(side="left", padx=2, expand=True)

        # Botón para generar CSV
        ttk.Button(
            botones_tree_frame,
            text="Generar CSV",
            command=self.generar_csv,
            style="Accent.TButton"
        ).pack(side="left", padx=2, expand=True)

        # Botón para generar reporte PPT
        ttk.Button(
            botones_tree_frame,
            text="Generar Reporte RCA",
            command=self.generar_reporte_ppt,
            style="Accent.TButton"
        ).pack(side="left", padx=2, expand=True)

    def generar_reporte_ppt(self):
        """Generar reporte de PowerPoint con los datos seleccionados"""
        seleccion = self.treeview.selection()
        
        if not seleccion:
            messagebox.showwarning("Advertencia", "Selecciona un registro primero")
            return

        # Obtener datos del registro seleccionado
        item = self.treeview.item(seleccion[0])
        lote = item["values"][0]
        
        # Buscar en los datos de Excel para obtener todos los campos
        datos = self.load_data()
        registro_completo = None
        for fila in datos:
            if str(fila[1]) == str(lote):
                registro_completo = {
                    'defect_code':" ",   # Código de defecto (placeholder)
                    'part_id': fila[2],  # Part ID
                    'lot_id': fila[1],   # Lote
                    'equipo': fila[3],   # Equipo
                    'comentario': fila[6],  # Comentario
                    'fecha': fila[8],    # Fecha
                    'turno': fila[10],   # Turno
                    'estatus': fila[12], # Estatus
                    'rca': fila[11],     # RCA
                    'hora': fila[9]      # Hora
                    
                }
                break
        
        if not registro_completo:
            messagebox.showerror("Error", "No se encontraron los datos completos del registro")
            return

        # Ejecutar en la misma aplicación
        report_window = tk.Toplevel(self.root)
        report_app = DefectReportApp(report_window, initial_data=registro_completo)

    def load_data(self):
        """Cargar datos desde el archivo Excel"""
        if not os.path.exists(self.ruta_base_datos):
            return []
        
        try:
            wb = openpyxl.load_workbook(self.ruta_base_datos)
            sheet = wb.active
            data = []
            
            # Comenzar desde la fila 2 (omitir encabezados)
            for row in sheet.iter_rows(min_row=2, values_only=True):
                # Filtrar filas vacías
                if any(cell is not None for cell in row):
                    data.append(row)
            
            return data
        except Exception as e:
            messagebox.showerror("Error", f"No se pudo cargar el archivo Excel:\n{str(e)}")
            return []

    def validar_lot(self, lot):
        """Validar formato del Lote (7 dígitos + . + 1 dígito)"""
        return re.match(r'^\d{7}\.\d{1}$', lot) is not None

    def guardar_datos(self, intentos=3):
        """Guardar los datos en el archivo Excel con confirmación y verificación de ID en tiempo real"""
        # Validar campos obligatorios
        if not self.lot_entry.get() or not self.part_id_entry.get() or not self.cometario_entry.get("1.0", "end-1c"):
            messagebox.showerror("Error", "¡Completa los campos obligatorios (Lote, Part ID, Comentario)!")
            return

        # Validar formato del Lote
        if not self.validar_lot(self.lot_entry.get()):
            messagebox.showerror("Error", "Formato de Lote inválido. Debe ser: 7 dígitos + . + 1 dígito")
            return

        # Validar hora (opcional)
        hora_val = self.hora_entry.get()
        if hora_val:
            try:
                datetime.strptime(hora_val, "%H:%M")
            except ValueError:
                messagebox.showerror("Error", "Formato de hora inválido. Usa HH:MM (24 hrs)")
                return

        # Mostrar resumen antes de guardar
        resumen = (
            f"Lote: {self.lot_entry.get()}\n"
            f"Part ID: {self.part_id_entry.get()}\n"
            f"Equipo: {self.equipo_combo.get()}\n"
            f"Comentario: {self.cometario_entry.get('1.0', 'end-1c')[:50]}...\n\n"
            "¿Deseas guardar este registro?"
        )
        
        confirmar = messagebox.askyesno("Confirmar Registro", resumen)
        if not confirmar:
            return

        for intento in range(intentos):
            try:
                if not os.path.exists(self.ruta_base_datos):
                    wb = Workbook()
                    ws = wb.active
                    ws.append([
                        "ID", "Lote", "Part ID", "Equipo", "CT", "CP", "Comentario",
                        "PCB", "Fecha", "Hora", "Turno", "RCA", "Estatus"
                    ])
                    wb.save(self.ruta_base_datos)
                    siguiente_id = 1  # Primer ID para archivo nuevo
                else:
                    # Cargar archivo y verificar ID justo antes de insertar
                    wb = openpyxl.load_workbook(self.ruta_base_datos)
                    ws = wb.active
                    
                    # Buscar el máximo ID y verificar duplicados
                    ids_existentes = set()
                    max_id = 0
                    for row in ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=1, max_col=1):
                        if row[0].value is not None:
                            try:
                                current_id = int(row[0].value)
                                ids_existentes.add(current_id)
                                if current_id > max_id:
                                    max_id = current_id
                            except (ValueError, TypeError):
                                continue
                    
                    siguiente_id = max_id + 1
                    
                    # Verificar si el ID ya existe (protección adicional)
                    while siguiente_id in ids_existentes:
                        siguiente_id += 1

                # Preparar datos para insertar
                datos_nuevos = [
                    siguiente_id,
                    self.lot_entry.get(),
                    self.part_id_entry.get(),
                    self.equipo_combo.get(),
                    self.ct_entry.get(),
                    self.cp_entry.get(),
                    self.cometario_entry.get("1.0", "end-1c"),
                    self.pcb_entry.get(),
                    self.calendario.get_date().strftime("%m/%d/%Y"),
                    hora_val,
                    self.turno_combo.get(),
                    self.rca_combo.get(),
                    self.estatus_combo.get()
                ]

                # Insertar y guardar
                ws.append(datos_nuevos)
                
                try:
                    wb.save(self.ruta_base_datos)
                    
                    # Mensaje de confirmación después de guardar
                    mensaje_exito = (
                        f"Registro #{siguiente_id} guardado correctamente\n\n"
                        f"Lote: {self.lot_entry.get()}\n"
                        f"Part ID: {self.part_id_entry.get()}\n"
                        f"Equipo: {self.equipo_combo.get()}\n"
                        f"Fecha: {self.calendario.get_date().strftime('%d/%m/%Y')}"
                    )
                    messagebox.showinfo("Éxito", mensaje_exito)
                    
                    self.limpiar_campos()
                    self.actualizar_treeview()
                    return  # Salir si todo fue exitoso
                    
                except PermissionError as pe:
                    if intento < intentos - 1:
                        # Esperar un poco antes de reintentar
                        time.sleep(0.5)
                        continue
                    raise  # Relanzar la excepción si se agotan los intentos

            except Exception as e:
                if intento == intentos - 1:  # Último intento fallido
                    error_msg = f"No se pudo guardar después de {intentos} intentos:\n{str(e)}"
                    if isinstance(e, PermissionError):
                        error_msg += "\n\nEl archivo está siendo usado por otro programa o usuario."
                    messagebox.showerror("Error", error_msg)
                continue

    def limpiar_campos(self):
        """Limpiar todos los campos del formulario"""
        respuesta = messagebox.askyesno("Confirmar Limpieza", "¿Estás seguro que deseas limpiar todos los campos?")
    
        if respuesta:
            self.lot_entry.delete(0, tk.END)
            self.lot_entry.insert(0, "Lote")
            self.part_id_entry.delete(0, tk.END)
            self.part_id_entry.insert(0, "Part ID")
            self.equipo_combo.set("Equipo")
            self.ct_entry.delete(0, tk.END)
            self.ct_entry.insert(0, "CT")
            self.cp_entry.delete(0, tk.END)
            self.cp_entry.insert(0, "CP")
            self.pcb_entry.delete(0, tk.END)
            self.pcb_entry.insert(0, "PCB")
            self.cometario_entry.delete("1.0", tk.END)
            self.turno_combo.set("Turno")
            self.rca_combo.set("RCA")
            self.estatus_combo.set("Estatus")
            self.hora_entry.delete(0, tk.END)
            self.id_a_modificar = None
            self.guardar_button.pack_forget()
            self.insert_button.pack(side="left", padx=2, expand=True)


    def actualizar_treeview(self):
        """Actualizar el Treeview con los datos del Excel"""
        for item in self.treeview.get_children():
            self.treeview.delete(item)
        
        try:
            data = self.load_data()
            for row in data:
                # Mostrar solo las columnas seleccionadas en el Treeview
                self.treeview.insert("", "end", values=(row[1], row[2], row[3], row[6], row[12]))
        except Exception as e:
            messagebox.showerror("Error", f"No se pudieron cargar los datos:\n{str(e)}")

    def copiar_seleccion(self):
        """Copiar los datos seleccionados al portapapeles con encabezados"""
        seleccion = self.treeview.selection()
        if not seleccion:
            messagebox.showwarning("Advertencia", "No hay datos seleccionados")
            return

        # Obtener los encabezados de las columnas
        encabezados = [self.treeview.heading(col)["text"] for col in self.treeview["columns"]]
        
        # Obtener los valores seleccionados
        item = self.treeview.item(seleccion[0])
        valores = item["values"]
        
        # Combinar encabezados y valores con formato
        texto_con_encabezados = ""
        for encabezado, valor in zip(encabezados, valores):
            texto_con_encabezados += f"{encabezado}: {valor}\n"
        
        # Copiar al portapapeles
        pyperclip.copy(texto_con_encabezados.strip())
        messagebox.showinfo("Éxito", "Datos copiados al portapapeles (con encabezados)")

    def modificar_evento(self):
        """Cargar los datos del registro seleccionado en el formulario para editar"""
        seleccion = self.treeview.selection()
        
        if not seleccion:
            messagebox.showwarning("Advertencia", "Selecciona un registro para modificar")
            return

        item = self.treeview.item(seleccion[0])
        valores = item["values"]
        lote = valores[0]

        # Cargar todos los datos desde Excel
        datos = self.load_data()
        for fila in datos:
            if str(fila[1]) == str(lote):
                self.id_a_modificar = fila[0]  # Guardamos el ID del registro
                
                # Actualizar campos del formulario
                self.lot_entry.delete(0, tk.END)
                self.lot_entry.insert(0, str(fila[1]))

                self.part_id_entry.delete(0, tk.END)
                self.part_id_entry.insert(0, str(fila[2]))

                self.equipo_combo.set(fila[3])
                self.ct_entry.delete(0, tk.END)
                self.ct_entry.insert(0, str(fila[4]))

                self.cp_entry.delete(0, tk.END)
                self.cp_entry.insert(0, str(fila[5]))

                self.cometario_entry.delete("1.0", tk.END)
                self.cometario_entry.insert("1.0", str(fila[6]))

                self.pcb_entry.delete(0, tk.END)
                self.pcb_entry.insert(0, str(fila[7]))

                try:
                    self.calendario.set_date(fila[8])
                except:
                    pass

                self.hora_entry.delete(0, tk.END)
                if fila[9] is not None:
                    self.hora_entry.insert(0, str(fila[9]))

                self.turno_combo.set(fila[10])
                self.rca_combo.set(fila[11])
                self.estatus_combo.set(fila[12])

                # Ocultar botón Insertar y mostrar Guardar Cambios
                self.insert_button.pack_forget()
                self.guardar_button.pack(side="left", padx=2, expand=True)
                return

        messagebox.showerror("Error", "No se pudo encontrar el registro completo en el Excel")

    def guardar_cambios(self):
        """Actualizar un registro existente en el Excel"""
        if self.id_a_modificar is None:
            messagebox.showerror("Error", "No hay registro en modo edición")
            return

        try:
            wb = openpyxl.load_workbook(self.ruta_base_datos)
            ws = wb.active

            for fila in ws.iter_rows(min_row=2):
                if fila[0].value == self.id_a_modificar:
                    fila[1].value = self.lot_entry.get()
                    fila[2].value = self.part_id_entry.get()
                    fila[3].value = self.equipo_combo.get()
                    fila[4].value = self.ct_entry.get()
                    fila[5].value = self.cp_entry.get()
                    fila[6].value = self.cometario_entry.get("1.0", "end-1c")
                    fila[7].value = self.pcb_entry.get()
                    fila[8].value = self.calendario.get_date().strftime("%m/%d/%Y")
                    fila[9].value = self.hora_entry.get()
                    fila[10].value = self.turno_combo.get()
                    fila[11].value = self.rca_combo.get()
                    fila[12].value = self.estatus_combo.get()
                    break

            wb.save(self.ruta_base_datos)
            messagebox.showinfo("Éxito", f"Registro #{self.id_a_modificar} actualizado correctamente")

            # Resetear estado
            self.id_a_modificar = None
            self.limpiar_campos()
            self.actualizar_treeview()
            self.guardar_button.pack_forget()
            self.insert_button.pack(side="left", padx=2, expand=True)
        except Exception as e:
            messagebox.showerror("Error", f"No se pudo guardar la modificación:\n{str(e)}")

    def buscar_por_lote(self):
        """Buscar registros por número de lote"""
        lote_buscado = self.lot_entry.get().strip()
        if not lote_buscado or lote_buscado == "Lote":
            messagebox.showwarning("Advertencia", "Ingresa un Lote para buscar")
            return
        
        for item in self.treeview.get_children():
            if self.treeview.item(item)["values"][0] == lote_buscado:
                self.treeview.selection_set(item)
                self.treeview.focus(item)
                self.treeview.see(item)
                return
        
        messagebox.showinfo("Información", f"No se encontró el Lote: {lote_buscado}")

    def generar_csv(self):
        """Generar un archivo CSV con todos los datos de la bitácora"""
        try:
            # Obtener los datos del Excel
            datos = self.load_data()
            
            if not datos:
                messagebox.showwarning("Advertencia", "No hay datos para exportar")
                return
            
            # Crear el nombre del archivo con fecha y hora
            fecha_hora = datetime.now().strftime("%Y%m%d_%H%M%S")
            nombre_archivo = f"Bitacora_CT_{fecha_hora}.csv"
            
            # Preguntar al usuario dónde guardar el archivo
            from tkinter import filedialog
            ruta_guardado = filedialog.asksaveasfilename(
                defaultextension=".csv",
                filetypes=[("Archivos CSV", "*.csv"), ("Todos los archivos", "*.*")],
                initialfile=nombre_archivo,
                title="Guardar archivo CSV como..."
            )
            
            if not ruta_guardado:  # El usuario canceló
                return
            
            # Escribir los datos en el archivo CSV
            import csv
            with open(ruta_guardado, mode='w', newline='', encoding='utf-8') as archivo:
                escritor = csv.writer(archivo)
                
                # Escribir encabezados
                encabezados = [
                    "ID", "Lote", "Part ID", "Equipo", "CT", "CP", "Comentario",
                    "PCB", "Fecha", "Hora", "Turno", "RCA", "Estatus"
                ]
                escritor.writerow(encabezados)
                
                # Escribir datos
                for fila in datos:
                    escritor.writerow(fila)
            
            messagebox.showinfo("Éxito", f"Archivo CSV generado correctamente:\n{ruta_guardado}")
            
        except Exception as e:
            messagebox.showerror("Error", f"No se pudo generar el archivo CSV:\n{str(e)}")

    def generar_grafica_turnos(self):
        """Generar gráfica de frecuencia de turnos por mes"""
        try:
            # Crear ventana emergente
            ventana_grafica = tk.Toplevel(self.root)
            ventana_grafica.title("Gráfica de Turnos por Mes")
            ventana_grafica.geometry("800x600")
            
            # Frame para controles
            control_frame = ttk.Frame(ventana_grafica)
            control_frame.pack(pady=10, fill="x", padx=20)
            
            # Etiqueta y combobox para selección de mes
            ttk.Label(control_frame, text="Seleccionar Mes:").grid(row=0, column=0, padx=5)
            
            meses = ["Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio", 
                    "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"]
            mes_combo = ttk.Combobox(control_frame, values=meses, width=15)
            mes_combo.grid(row=0, column=1, padx=5)
            mes_combo.current(0)  # Seleccionar enero por defecto
            
            # Etiqueta y combobox para selección de año
            ttk.Label(control_frame, text="Seleccionar Año:").grid(row=0, column=2, padx=5)
            
            # Obtener el año actual
            año_actual = datetime.now().year
            años = [str(año_actual - 1), str(año_actual), str(año_actual + 1)]
            año_combo = ttk.Combobox(control_frame, values=años, width=8)
            año_combo.grid(row=0, column=3, padx=5)
            año_combo.set(str(año_actual))  # Seleccionar año actual por defecto
            
            # Botón para generar gráfica
            btn_generar = ttk.Button(control_frame, text="Generar Gráfica", 
                                    command=lambda: self.mostrar_grafica_turnos(mes_combo.get(), año_combo.get(), ventana_grafica),
                                    style="Accent.TButton")
            btn_generar.grid(row=0, column=4, padx=10)
            
            # Frame para la gráfica
            grafica_frame = ttk.Frame(ventana_grafica)
            grafica_frame.pack(fill="both", expand=True, padx=20, pady=10)
            
        except Exception as e:
            messagebox.showerror("Error", f"No se pudo crear la ventana de gráficas:\n{str(e)}")

    def mostrar_grafica_turnos(self, mes_seleccionado, año_seleccionado, ventana):
        """Mostrar gráfica de frecuencia de turnos para el mes seleccionado"""
        try:
            # Obtener datos del Excel
            data = self.load_data()
            if not data:
                messagebox.showinfo("Información", "No hay datos para mostrar")
                return
            
            # Obtener número de mes (1-12)
            meses = ["Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio", 
                     "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"]
            num_mes = meses.index(mes_seleccionado) + 1

            # Inicializar contadores
            turnos = {"A": 0, "B": 0, "C": 0, "D": 0}
            total_eventos = 0
            
            for fila in data:
                try:
                    fecha = fila[8]  # Columna de Fecha

                    # Si es cadena, convertirla a datetime
                    if isinstance(fecha, str):
                        fecha = datetime.strptime(fecha, "%m/%d/%Y")

                    # Verificar mes y año
                    if fecha.month == num_mes and fecha.year == int(año_seleccionado):
                        turno = str(fila[10]).strip().upper()  # Normalizar turno
                        if turno in turnos:
                            turnos[turno] += 1
                            total_eventos += 1
                except Exception as e:
                    print(f"Error al procesar fila: {fila}\n{e}")
                    continue

            if total_eventos == 0:
                messagebox.showinfo("Información", f"No hay eventos para {mes_seleccionado} {año_seleccionado}")
                return

            # Crear figura
            fig, ax = plt.subplots(figsize=(8, 5))
            turnos_keys = list(turnos.keys())
            valores = list(turnos.values())

            bars = ax.bar(turnos_keys, valores, color=['#4e79a7', '#f28e2c', '#e15759', '#76b7b2'])

            for bar in bars:
                height = bar.get_height()
                ax.annotate(f'{height}',
                            xy=(bar.get_x() + bar.get_width() / 2, height),
                            xytext=(0, 3),
                            textcoords="offset points",
                            ha='center', va='bottom')

            ax.set_title(f'Frecuencia de Turnos - {mes_seleccionado} {año_seleccionado}', fontsize=14)
            ax.set_xlabel('Turno', fontsize=12)
            ax.set_ylabel('Cantidad de Eventos', fontsize=12)
            ax.set_ylim(0, max(valores) * 1.2)
            ax.grid(axis='y', linestyle='--', alpha=0.7)
            plt.tight_layout()

            # Limpiar frame anterior
            for widget in ventana.winfo_children():
                if isinstance(widget, ttk.Frame) and widget.winfo_name() == "!frame2":
                    widget.destroy()

            grafica_frame = ttk.Frame(ventana, name="frame2")
            grafica_frame.pack(fill="both", expand=True, padx=20, pady=10)

            canvas = FigureCanvasTkAgg(fig, master=grafica_frame)
            canvas.draw()
            canvas.get_tk_widget().pack(fill="both", expand=True)

        except Exception as e:
            messagebox.showerror("Error", f"No se pudo generar la gráfica:\n{str(e)}")

    def generar_grafica_equipos(self):
        """Ventana para seleccionar mes y año para la gráfica por equipos"""
        try:
            ventana_equipos = tk.Toplevel(self.root)
            ventana_equipos.title("Gráfica por Equipos")
            ventana_equipos.geometry("900x600")

            control_frame = ttk.Frame(ventana_equipos)
            control_frame.pack(pady=10, fill="x", padx=20)

            # Selección de mes
            ttk.Label(control_frame, text="Mes:").grid(row=0, column=0, padx=5)
            meses = ["Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio", 
                     "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"]
            mes_combo = ttk.Combobox(control_frame, values=meses, width=15)
            mes_combo.grid(row=0, column=1, padx=5)
            mes_combo.current(0)

            # Selección de año
            ttk.Label(control_frame, text="Año:").grid(row=0, column=2, padx=5)
            año_actual = datetime.now().year
            años = [str(año_actual - 1), str(año_actual), str(año_actual + 1)]
            año_combo = ttk.Combobox(control_frame, values=años, width=8)
            año_combo.grid(row=0, column=3, padx=5)
            año_combo.set(str(año_actual))

            # Botón generar
            btn_generar = ttk.Button(control_frame, text="Generar Gráfica",
                                     command=lambda: self.mostrar_grafica_equipos(mes_combo.get(), año_combo.get(), ventana_equipos),
                                     style="Accent.TButton")
            btn_generar.grid(row=0, column=4, padx=10)

            # Frame para la gráfica
            grafica_frame = ttk.Frame(ventana_equipos)
            grafica_frame.pack(fill="both", expand=True, padx=20, pady=10)

        except Exception as e:
            messagebox.showerror("Error", f"No se pudo crear la ventana de gráficas:\n{str(e)}")

    def mostrar_grafica_equipos(self, mes_seleccionado, año_seleccionado, ventana):
        """Genera gráfica de eventos por equipo, separados por turno"""
        try:
            data = self.load_data()
            if not data:
                messagebox.showinfo("Información", "No hay datos para mostrar")
                return

            meses = ["Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio", 
                     "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"]
            num_mes = meses.index(mes_seleccionado) + 1

            # Diccionario: equipo -> {turno -> cantidad}
            equipos_data = {}

            for fila in data:
                try:
                    fecha = fila[8]
                    if isinstance(fecha, str):
                        fecha = datetime.strptime(fecha, "%m/%d/%Y")
                    if fecha.month != num_mes or fecha.year != int(año_seleccionado):
                        continue

                    equipo = str(fila[3]).strip()
                    turno = str(fila[10]).strip().upper()

                    if equipo and turno in ["A", "B", "C", "D"]:
                        if equipo not in equipos_data:
                            equipos_data[equipo] = {"A": 0, "B": 0, "C": 0, "D": 0}
                        equipos_data[equipo][turno] += 1

                except Exception as e:
                    print(f"Error al procesar fila: {fila}\n{e}")
                    continue

            if not equipos_data:
                messagebox.showinfo("Información", f"No hay eventos para {mes_seleccionado} {año_seleccionado}")
                return

            # Preparar datos para gráfica
            equipos = sorted(equipos_data.keys())
            turnos = ["A", "B", "C", "D"]
            colores = ['#4e79a7', '#f28e2c', '#e15759', '#76b7b2']

            valores_turnos = {t: [equipos_data[e].get(t, 0) for e in equipos] for t in turnos}

            fig, ax = plt.subplots(figsize=(10, 6))

            bottom = [0] * len(equipos)
            for i, turno in enumerate(turnos):
                ax.bar(equipos, valores_turnos[turno], bottom=bottom, label=f'Turno {turno}', color=colores[i])
                bottom = [sum(x) for x in zip(bottom, valores_turnos[turno])]

            ax.set_title(f'Eventos por Equipo y Turno - {mes_seleccionado} {año_seleccionado}', fontsize=14)
            ax.set_xlabel("Equipo", fontsize=12)
            ax.set_ylabel("Cantidad de Eventos", fontsize=12)
            ax.legend()
            ax.set_xticklabels(equipos, rotation=45, ha='right')
            ax.grid(axis='y', linestyle='--', alpha=0.5)

            plt.tight_layout()

            # Mostrar en ventana
            grafica_frame = ttk.Frame(ventana)
            grafica_frame.pack(fill="both", expand=True, padx=20, pady=10)

            canvas = FigureCanvasTkAgg(fig, master=grafica_frame)
            canvas.draw()
            canvas.get_tk_widget().pack(fill="both", expand=True)

        except Exception as e:
            messagebox.showerror("Error", f"No se pudo generar la gráfica:\n{str(e)}")

    def centrar_ventana(self):
        self.root.update()
        self.root.minsize(1300, 650)
        x = (self.root.winfo_screenwidth() // 2) - (self.root.winfo_width() // 2)
        y = (self.root.winfo_screenheight() // 2) - (self.root.winfo_height() // 2)
        self.root.geometry(f"+{x}+{y}")

####################################################################################################

class DefectReportApp:
    def __init__(self, root, initial_data=None):
        self.root = root
        self.root.title("Sistema de Reporte de Defectos")
        self.root.geometry("800x800")
        
        # Configurar el tema oscuro
        self.configurar_tema_oscuro()
        
        # Variables para almacenar los datos
        self.part_id = tk.StringVar()
        self.lot_id = tk.StringVar()
        self.defect_code = tk.StringVar()
        self.date = tk.StringVar()
        self.shift = tk.StringVar()
        self.what = tk.StringVar()
        self.why = tk.StringVar()
        self.where = tk.StringVar()
        self.when = tk.StringVar()
        self.who = tk.StringVar()
        self.how = tk.StringVar()
        self.how_much = tk.StringVar()
        self.occurrence = tk.StringVar()
        self.detection = tk.StringVar()
        self.systemic = tk.StringVar()
        
        # Referencias a los widgets Text
        self.how_entry = None
        self.occurrence_entry = None
        self.detection_entry = None
        self.systemic_entry = None
        
        # Crear interfaz
        self.create_widgets()
        
        # Cargar datos iniciales si se proporcionan
        if initial_data:
            self.load_initial_data(initial_data)

    def configurar_tema_oscuro(self):
        """Configurar el tema forest-dark solo si no está ya cargado"""
        try:
            # Revisar si el tema ya existe
            if "forest-dark" not in self.root.tk.call("ttk::themes"):
                self.root.tk.call("source", recurso_relativo("forest-dark.tcl"))
            
            style = ttk.Style(self.root)
            style.theme_use("forest-dark")

            # Configuración adicional de colores
            style.configure('TFrame',)
            style.configure('TLabel', foreground='white', font=('Arial', 10))
            style.configure('TButton', font=('Arial', 10))
            style.configure('Header.TLabel', font=('Arial', 12, 'bold'), foreground='white')
            style.configure('TCombobox', fieldbackground='#2d2d2d', foreground='white')
            style.configure('TEntry', fieldbackground='#2d2d2d', foreground='white')
            
            self.root.configure(bg='#1e1e1e')

        except Exception as e:
            messagebox.showerror("Error", f"No se pudo cargar el tema: {str(e)}")

    def create_widgets(self):
        # Frame principal
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # Frame para datos básicos
        basic_frame = ttk.LabelFrame(main_frame, text="Datos Básicos del Defecto", padding="10")
        basic_frame.pack(fill=tk.X, pady=5)
        
        # Part ID
        ttk.Label(basic_frame, text="Part ID:*").grid(row=0, column=0, sticky=tk.W, padx=5, pady=2)
        ttk.Entry(basic_frame, textvariable=self.part_id, width=30).grid(row=0, column=1, sticky=tk.W, padx=5, pady=2)
        
        # Lot ID
        ttk.Label(basic_frame, text="Lot ID:*").grid(row=1, column=0, sticky=tk.W, padx=5, pady=2)
        ttk.Entry(basic_frame, textvariable=self.lot_id, width=30).grid(row=1, column=1, sticky=tk.W, padx=5, pady=2)
        
        # Código Defecto
        ttk.Label(basic_frame, text="Código Defecto:*").grid(row=2, column=0, sticky=tk.W, padx=5, pady=2)
        defect_codes = ['KI', 'MH', 'KQ', 'KF', 'KX', "Otro"]
        ttk.Combobox(basic_frame, textvariable=self.defect_code, values=defect_codes, width=27).grid(row=2, column=1, sticky=tk.W, padx=5, pady=2)
        
        # Fecha
        ttk.Label(basic_frame, text="Fecha:*").grid(row=3, column=0, sticky=tk.W, padx=5, pady=2)
        self.date_entry = DateEntry(basic_frame, textvariable=self.date, date_pattern='dd/mm/yyyy', width=27)
        self.date_entry.grid(row=3, column=1, sticky=tk.W, padx=5, pady=2)
        
        # Turno
        ttk.Label(basic_frame, text="Turno:*").grid(row=4, column=0, sticky=tk.W, padx=5, pady=2)
        ttk.Combobox(basic_frame, textvariable=self.shift, values=["Turno A", "Turno B", "Turno C", "Turno D"], width=27).grid(row=4, column=1, sticky=tk.W, padx=5, pady=2)
        
        # Frame para análisis 5W2H
        analysis_frame = ttk.LabelFrame(main_frame, text="Análisis del Defecto (5W2H)", padding="10")
        analysis_frame.pack(fill=tk.BOTH, expand=True, pady=5)
        
        # Configurar el fondo oscuro para los widgets Text
        text_bg = '#2d2d2d'
        text_fg = 'white'
        
        # What
        ttk.Label(analysis_frame, text="What (Que ocurrio):*").grid(row=0, column=0, sticky=tk.W, padx=5, pady=2)
        ttk.Entry(analysis_frame, textvariable=self.what, width=60).grid(row=0, column=1, sticky=tk.W, padx=5, pady=2)
        
        # Why
        ttk.Label(analysis_frame, text="Why (Por que ocurrio):*").grid(row=1, column=0, sticky=tk.W, padx=5, pady=2)
        ttk.Entry(analysis_frame, textvariable=self.why, width=60).grid(row=1, column=1, sticky=tk.W, padx=5, pady=2)
        
        # Where
        ttk.Label(analysis_frame, text="Where (Dónde se detectó):*").grid(row=2, column=0, sticky=tk.W, padx=5, pady=2)
        ttk.Entry(analysis_frame, textvariable=self.where, width=60).grid(row=2, column=1, sticky=tk.W, padx=5, pady=2)
        
        # Who
        ttk.Label(analysis_frame, text="Who (Quién lo detectó):*").grid(row=3, column=0, sticky=tk.W, padx=5, pady=2)
        ttk.Entry(analysis_frame, textvariable=self.who, width=60).grid(row=3, column=1, sticky=tk.W, padx=5, pady=2)
        
        # How
        ttk.Label(analysis_frame, text="How (Cómo se detectó):*").grid(row=4, column=0, sticky=tk.NW, padx=5, pady=2)
        self.how_entry = tk.Text(
            analysis_frame, 
            width=60, 
            height=3, 
            wrap=tk.WORD,
            bg=text_bg,
            fg=text_fg,
            insertbackground='white'
        )
        self.how_entry.grid(row=4, column=1, padx=5, pady=2)
        self.how_entry.bind('<KeyRelease>', lambda e: self.how.set(self.how_entry.get("1.0", tk.END)))
        
        # How much
        ttk.Label(analysis_frame, text="How much (Impacto/Cantidad):*").grid(row=5, column=0, sticky=tk.W, padx=5, pady=2)
        ttk.Entry(analysis_frame, textvariable=self.how_much, width=60).grid(row=5, column=1, sticky=tk.W, padx=5, pady=2)
        
        # Frame para Posible Causa Raíz
        root_cause_frame = ttk.LabelFrame(main_frame, text="Posible Causa Raíz", padding="10")
        root_cause_frame.pack(fill=tk.BOTH, pady=5)
        
        # Occurrence
        ttk.Label(root_cause_frame, text="Occurrence:").grid(row=0, column=0, sticky=tk.W, padx=5, pady=2)
        self.occurrence_entry = tk.Text(
            root_cause_frame, 
            width=60, 
            height=3, 
            wrap=tk.WORD,
            bg=text_bg,
            fg=text_fg,
            insertbackground='white'
        )
        self.occurrence_entry.grid(row=0, column=1, padx=5, pady=2)
        self.occurrence_entry.bind('<KeyRelease>', lambda e: self.occurrence.set(self.occurrence_entry.get("1.0", tk.END)))
        
        # Detection
        ttk.Label(root_cause_frame, text="Detection:").grid(row=1, column=0, sticky=tk.W, padx=5, pady=2)
        self.detection_entry = tk.Text(
            root_cause_frame, 
            width=60, 
            height=3, 
            wrap=tk.WORD,
            bg=text_bg,
            fg=text_fg,
            insertbackground='white'
        )
        self.detection_entry.grid(row=1, column=1, padx=5, pady=2)
        self.detection_entry.bind('<KeyRelease>', lambda e: self.detection.set(self.detection_entry.get("1.0", tk.END)))
        
        # Systemic
        ttk.Label(root_cause_frame, text="Systemic:").grid(row=2, column=0, sticky=tk.W, padx=5, pady=2)
        self.systemic_entry = tk.Text(
            root_cause_frame, 
            width=60, 
            height=3, 
            wrap=tk.WORD,
            bg=text_bg,
            fg=text_fg,
            insertbackground='white'
        )
        self.systemic_entry.grid(row=2, column=1, padx=5, pady=2)
        self.systemic_entry.bind('<KeyRelease>', lambda e: self.systemic.set(self.systemic_entry.get("1.0", tk.END)))
        
        # Frame para botones
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill=tk.X, pady=10)
        
        # Botones
        ttk.Button(
            button_frame, 
            text="Generar Reporte", 
            command=self.generate_report,
            style="Accent.TButton"
        ).pack(side=tk.LEFT, padx=5)
        
        ttk.Button(
            button_frame, 
            text="Limpiar Formulario", 
            command=self.clear_form
        ).pack(side=tk.LEFT, padx=5)
        
        ttk.Button(
            button_frame, 
            text="Salir", 
            command=self.root.quit
        ).pack(side=tk.RIGHT, padx=5)
    
    def load_initial_data(self, data):
        """Cargar datos iniciales en el formulario"""
        self.part_id.set(data.get('part_id', ''))
        self.lot_id.set(data.get('lot_id', ''))
        self.defect_code.set(data.get('defect_code', ''))
        
        # Convertir fecha de formato mm/dd/yyyy a dd/mm/yyyy si es necesario
        fecha = data.get('fecha', '')
        if fecha:
            try:
                if '/' in fecha:
                    mes, dia, anio = fecha.split('/')
                    self.date.set(f"{dia}/{mes}/{anio}")
                else:
                    self.date.set(fecha)
            except:
                self.date.set(fecha)
        
        turno = data.get('turno', '')
        if turno:
            self.shift.set(f"Turno {turno}" if len(turno) == 1 else turno)
        
        # Autocompletar campos de análisis con el comentario
        comentario = data.get('comentario', '')
        if comentario:
            self.what.set("Descripción del defecto")
            self.why.set("Que fue lo que causó el defecto")
            self.where.set(f"Equipo {data.get('equipo', '')}")
            self.who.set("Operador/Inspector")
            self.how_entry.insert("1.0", "Descripción corta de que fue lo que ocurrió")
            self.how_much.set("Pendiente de determinar")
            
            # Sugerir posibles causas basadas en el equipo
            equipo = data.get('equipo', '')
            if equipo.startswith('Towa'):
                self.occurrence_entry.insert("1.0", comentario)
                self.detection_entry.insert("1.0", "Cuando el equipo Towa detecta un defecto, se detiene automáticamente")
                self.systemic_entry.insert("1.0", "El equipo No cuenta con un sistema de monitoreo adecuado")
    
    def validate_fields(self):
        required_fields = [
            (self.part_id, "Part ID"),
            (self.lot_id, "Lot ID"),
            (self.defect_code, "Código Defecto"),
            (self.date, "Fecha"),
            (self.shift, "Turno"),
            (self.what, "What"),
            (self.why, "Why"),
            (self.where, "Where"),
            (self.who, "Who"),
            (self.how, "How"),
            (self.how_much, "How much")
        ]
        
        for field, name in required_fields:
            if not field.get().strip():
                messagebox.showerror("Error", f"El campo {name} es obligatorio")
                return False
    
        return True
    
    def clear_form(self):
        """Limpiar todos los campos del formulario"""
        # Limpiar todas las variables
        for var in [self.part_id, self.lot_id, self.defect_code, self.date, 
                   self.shift, self.what, self.why, self.where, 
                   self.when, self.who, self.how, self.how_much,
                   self.occurrence, self.detection, self.systemic]:
            var.set("")
        
        # Limpiar los widgets Text manualmente
        self.how_entry.delete("1.0", tk.END)
        self.occurrence_entry.delete("1.0", tk.END)
        self.detection_entry.delete("1.0", tk.END)
        self.systemic_entry.delete("1.0", tk.END)
        
        # Restablecer valores por defecto
        self.date.set(datetime.now().strftime('%d/%m/%Y'))
        self.shift.set("Turno A")
    
    def generate_report(self):
        """Generar el reporte de PowerPoint"""
        if not self.validate_fields():
            return
            
        try:
            replacements = {
                '%PartID%': self.part_id.get(),
                '%LotID%': self.lot_id.get(),
                '%CodigoDefecto%': self.defect_code.get(),
                '%Fecha%': self.date.get(),
                '%Turno%': self.shift.get(),
                '%What%': self.what.get(),
                '%Why%': self.why.get(),
                '%Where%': self.where.get(),
                '%When%': self.when.get(),
                '%Who%': self.who.get(),
                '%How%': self.how.get(),
                '%HowMuch%': self.how_much.get(),
                '%Occurrence%': self.occurrence.get(),
                '%Detection%': self.detection.get(),
                '%Systemic%': self.systemic.get()
            }
            
            # Buscar la plantilla en varias ubicaciones posibles
            template_paths = [
                recurso_relativo("ppt_test.pptx"),
                recurso_relativo(os.path.join("templates", "ppt_test.pptx")),
                "ppt_test.pptx",
                os.path.join(os.path.dirname(__file__), "templates", "ppt_test.pptx")
            ]
            
            template_path = None
            for path in template_paths:
                if os.path.exists(path):
                    template_path = path
                    break
            
            if not template_path:
                messagebox.showerror("Error", "No se encontró la plantilla ppt_test.pptx")
                return
                
            prs = Presentation(template_path)
            
            def replace_in_shape(shape):
                if hasattr(shape, 'text'):
                    original_text = shape.text
                    for key, value in replacements.items():
                        if key in original_text:
                            shape.text = original_text.replace(key, str(value))
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    font = run.font
                                    font.name = 'Calibri'
                                    font.size = Pt(12)
                                    font.bold = False
                                    font.italic = True
                
                elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            original_text = run.text
                            for key, value in replacements.items():
                                if key in original_text:
                                    run.text = original_text.replace(key, str(value))
                                    font = run.font
                                    font.name = 'Calibri (Body)'
                                    font.size = Pt(12)
                                    font.bold = False
                                    font.italic = True
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    replace_in_shape(shape)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                replace_in_shape(cell)
                    if shape.shape_type == 6:  # Group shape
                        for subshape in shape.shapes:
                            replace_in_shape(subshape)
            
            # Mostrar diálogo para guardar
            default_filename = f"Reporte_{self.part_id.get()}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            file_path = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")],
                initialfile=default_filename,
                title="Guardar reporte como..."
            )
            
            if not file_path:
                return
            
            prs.save(file_path)
            
            messagebox.showinfo("Éxito", f"Reporte generado exitosamente en:\n{file_path}")
            
            # Preguntar si desea abrir el archivo
            if messagebox.askyesno("Abrir archivo", "¿Desea abrir el reporte generado?"):
                if os.name == 'nt':
                    os.startfile(file_path)
                else:
                    opener = 'open' if sys.platform == 'darwin' else 'xdg-open'
                    subprocess.call([opener, file_path])
        
        except Exception as e:
            messagebox.showerror("Error", f"Ocurrió un error al generar el reporte:\n{str(e)}")

####################################################################################################
if __name__ == "__main__":
    root = tk.Tk()
    app = BitacoraCTApp(root)
    root.mainloop()
